!pip install -q google-generativeai PyMuPDF python-pptx pandas

import os
import google.generativeai as genai
import fitz  # PyMuPDF
from pptx import Presentation
import pandas as pd

# Configuration of API Key
GOOGLE_API_KEY = "AIzaSyC_4vQF76i-gCwWQE3l8ILVnskb6eqfPXA"
genai.configure(api_key=GOOGLE_API_KEY)
model = genai.GenerativeModel("gemini-2.0-flash-exp")


# Data Extraction From The Files
def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_excel(xls_path):
    text = ""
    xls = pd.read_excel(xls_path, sheet_name=None)
    for sheet_name, df in xls.items():
        text += f"\nSheet: {sheet_name}\n"
        text += df.to_string(index=False)
    return text


# Storing the Extracted Data in Variables
pdf_text = extract_text_from_pdf("school.pdf")
pptx_text = extract_text_from_pptx("school.pptx")
xls_text = extract_text_from_excel("school.xls")


# Combining the data Stored in Above Variables
full_context = pdf_text + "\n" + pptx_text + "\n" + xls_text


# User Interface to Ask Questions
def ask_question(user_query):
    prompt = f"""Use the following data from PDF, PPTX and Excel files to answer the question.

--- START OF CONTEXT ---
{full_context[:15000]}  # Limit context to fit model token limit
--- END OF CONTEXT ---

Question: {user_query}
Answer:"""

    response = model.generate_content(prompt)
    return response.text


# Keepint the question Interface For Multiple Queris with Single Code Run

while True:
    question = input("Ask a question (or type 'exit' to quit): ").strip()
    if question.lower() == "exit":
        print("Thanks for using the RAG chatbot!")
        break
    answer = ask_question(question)
    print("Answer:", answer)
    print("\n" + "-"*50 + "\n")

# the end